"""Multi-MIME text extraction for the Vault file-upload feature.

Single ``extract_text(content, mime_type)`` entry point dispatches to a
per-format extractor and returns plain UTF-8 text. The output text is
the corpus that downstream chunking + embedding consume — quality of
extraction directly drives RAG quality, so we keep the per-format
helpers small and well-tested rather than wrapping a heavyweight
extractor like Apache Tika.

Supported (matches plans/vault-rag-2026-05-04.md Section 2):
    PDF, DOCX, PPTX, XLSX, TXT, MD, RTF, CSV, HTML, JSON

Each MIME's library is imported lazily so a single missing optional
dep doesn't break the test suite or unrelated code paths.
"""

from __future__ import annotations

import csv
import io
import json
import logging
import re
from typing import Callable

logger = logging.getLogger(__name__)


SUPPORTED_MIME_TYPES: frozenset[str] = frozenset(
    {
        "application/pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "text/plain",
        "text/markdown",
        "text/html",
        "text/csv",
        "application/rtf",
        "text/rtf",
        "application/json",
    }
)


class VaultExtractionError(RuntimeError):
    """Raised when bytes can't be turned into text."""


def is_supported(mime_type: str) -> bool:
    """Cheap check the endpoint runs before accepting an upload."""
    return mime_type.lower() in SUPPORTED_MIME_TYPES


def _extract_pdf(content: bytes) -> str:
    """PDF via pdfplumber. Reuses the same library the existing
    ``pdf_processor`` path uses for FOCUS reports — no new dep."""
    import pdfplumber

    with pdfplumber.open(io.BytesIO(content)) as pdf:
        pages = []
        for page in pdf.pages:
            page_text = page.extract_text() or ""
            if page_text.strip():
                pages.append(page_text)
        return "\n\n".join(pages)


def _extract_docx(content: bytes) -> str:
    from docx import Document

    document = Document(io.BytesIO(content))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
    # Tables aren't covered by paragraph iteration; flatten cells too.
    for table in document.tables:
        for row in table.rows:
            row_text = " | ".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                parts.append(row_text)
    return "\n".join(parts)


def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(content))
    slides: list[str] = []
    for index, slide in enumerate(presentation.slides):
        slide_lines = [f"[Slide {index + 1}]"]
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = "".join(run.text for run in paragraph.runs)
                if text.strip():
                    slide_lines.append(text)
        if len(slide_lines) > 1:
            slides.append("\n".join(slide_lines))
    return "\n\n".join(slides)


def _extract_xlsx(content: bytes) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(content), data_only=True, read_only=True)
    sheets: list[str] = []
    for sheet_name in workbook.sheetnames:
        worksheet = workbook[sheet_name]
        rows: list[str] = [f"[Sheet: {sheet_name}]"]
        for row in worksheet.iter_rows(values_only=True):
            cells = ["" if value is None else str(value) for value in row]
            if any(cell.strip() for cell in cells):
                rows.append(" | ".join(cells))
        if len(rows) > 1:
            sheets.append("\n".join(rows))
    workbook.close()
    return "\n\n".join(sheets)


def _extract_rtf(content: bytes) -> str:
    from striprtf.striprtf import rtf_to_text

    return rtf_to_text(content.decode("utf-8", errors="replace"))


def _decode_text(content: bytes) -> str:
    """UTF-8 with replacement for TXT/MD/CSV/JSON."""
    return content.decode("utf-8", errors="replace")


_HTML_SCRIPT_STYLE = re.compile(
    r"<(script|style)\b[^>]*>.*?</\1>", flags=re.IGNORECASE | re.DOTALL
)
_HTML_TAG = re.compile(r"<[^>]+>")
_HTML_ENTITY_TRIVIAL = {
    "&nbsp;": " ",
    "&amp;": "&",
    "&lt;": "<",
    "&gt;": ">",
    "&quot;": '"',
    "&apos;": "'",
}


def _extract_html(content: bytes) -> str:
    """Naive HTML-to-text. Not a full parser — strips <script>/<style>,
    drops remaining tags, decodes a handful of common entities. Good
    enough for RAG context where surrounding text matters more than
    semantic structure. If users start uploading raw web pages with
    heavy JS we'd need ``beautifulsoup4`` — keeping the dep surface
    minimal until that's a real complaint."""
    text = _decode_text(content)
    text = _HTML_SCRIPT_STYLE.sub(" ", text)
    text = _HTML_TAG.sub(" ", text)
    for entity, char in _HTML_ENTITY_TRIVIAL.items():
        text = text.replace(entity, char)
    # Collapse runs of whitespace introduced by tag stripping.
    return re.sub(r"\s+", " ", text).strip()


def _extract_csv(content: bytes) -> str:
    """Sniff delimiter, return as pipe-separated rows for chunkable text."""
    text = _decode_text(content)
    sniffer = csv.Sniffer()
    try:
        dialect = sniffer.sniff(text[:4096], delimiters=",;\t|")
    except csv.Error:
        dialect = csv.excel
    rows = list(csv.reader(io.StringIO(text), dialect))
    return "\n".join(" | ".join(cell.strip() for cell in row) for row in rows)


def _extract_json(content: bytes) -> str:
    """Pretty-print JSON for readability in the embedded chunks."""
    try:
        parsed = json.loads(_decode_text(content))
    except json.JSONDecodeError as exc:
        raise VaultExtractionError(f"Invalid JSON: {exc}") from exc
    return json.dumps(parsed, indent=2, ensure_ascii=False)


_DISPATCH: dict[str, Callable[[bytes], str]] = {
    "application/pdf": _extract_pdf,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": _extract_docx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": _extract_pptx,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": _extract_xlsx,
    "text/plain": _decode_text,
    "text/markdown": _decode_text,
    "text/html": _extract_html,
    "text/csv": _extract_csv,
    "application/rtf": _extract_rtf,
    "text/rtf": _extract_rtf,
    "application/json": _extract_json,
}


def extract_text(content: bytes, mime_type: str) -> str:
    """Dispatch by MIME, return cleaned plain text.

    Raises ``VaultExtractionError`` for unsupported MIME or any per-format
    parse failure. The endpoint maps that to a ``failed`` row so the FE
    can show a yellow chip + Retry button.
    """
    mime = mime_type.lower()
    extractor = _DISPATCH.get(mime)
    if extractor is None:
        raise VaultExtractionError(
            f"Unsupported MIME type for vault extraction: {mime_type!r}"
        )
    try:
        text = extractor(content)
    except VaultExtractionError:
        raise
    except Exception as exc:
        # Per-format libraries raise wide variety of exceptions; collapse
        # them into our domain error so the caller can switch on one type.
        raise VaultExtractionError(
            f"Failed to extract text from {mime_type} content: {exc}"
        ) from exc

    # Final guard: if the file was technically valid but yielded nothing
    # useful (image-only PDF, empty workbook), surface as failed so the
    # user sees the yellow chip instead of an empty file appearing
    # successfully embedded.
    if not text or not text.strip():
        raise VaultExtractionError(
            "Extracted text was empty — file may be image-only, "
            "password-protected, or otherwise unreadable."
        )
    return text
